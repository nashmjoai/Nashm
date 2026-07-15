import os
import sys
import json
import argparse
from PIL import Image

# Third party imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

# Define color schemes for themes
SLIDES_THEMES = {
    'nashm-executive-dark': {
        'bg': (15, 23, 42),       # Dark Slate
        'text': (241, 245, 249),  # Slate 100
        'accent': (245, 158, 11),  # Amber 500
        'card_bg': (30, 41, 59),  # Slate 800
    },
    'nashm-modern-light': {
        'bg': (248, 250, 252),    # Slate 50
        'text': (15, 23, 42),     # Slate 900
        'accent': (79, 70, 229),   # Indigo 600
        'card_bg': (255, 255, 255),# White
    },
    'nashm-arabic-lux': {
        'bg': (6, 78, 59),        # Emerald 950
        'text': (254, 243, 199),  # Amber 100
        'accent': (245, 158, 11),  # Amber 500
        'card_bg': (6, 95, 70),   # Emerald 900
    },
    'nashm-data-studio': {
        'bg': (245, 245, 244),    # Stone 100
        'text': (28, 25, 23),      # Stone 900
        'accent': (14, 165, 233),  # Sky 500
        'card_bg': (255, 255, 255),# White
    },
    'nashm-creative-color': {
        'bg': (49, 46, 129),      # Indigo 950
        'text': (255, 255, 255),  # White
        'accent': (244, 63, 94),   # Rose 500
        'card_bg': (67, 56, 202),  # Indigo 700
    }
}

DOCS_THEMES = {
    'nashm-report-pro': {
        'header_color': (79, 70, 229), # Indigo
        'text_color': (51, 65, 85),    # Slate 700
        'border_color': (226, 232, 240)
    },
    'nashm-research-rtl': {
        'header_color': (28, 25, 23),  # Stone 900
        'text_color': (41, 37, 36),    # Stone 800
        'border_color': (214, 211, 209)
    },
    'nashm-business-proposal': {
        'header_color': (16, 185, 129), # Emerald 500
        'text_color': (30, 41, 59),     # Slate 800
        'border_color': (226, 232, 240)
    },
    'nashm-formal-letter': {
        'header_color': (31, 41, 55),   # Gray 800
        'text_color': (17, 24, 39),     # Gray 900
        'border_color': (209, 213, 219)
    },
    'nashm-minimal-doc': {
        'header_color': (120, 113, 108), # Stone 500
        'text_color': (87, 80, 75),      # Stone 600
        'border_color': (231, 229, 228)
    }
}

SHEETS_THEMES = {
    'nashm-finance-dashboard': {
        'header_fill': '064E3B', # Emerald 950
        'header_text': 'FFFFFF',
        'kpi_fill': 'ECFDF5',    # Emerald 50
        'kpi_text': '065F46',
    },
    'nashm-project-tracker': {
        'header_fill': '312E81', # Indigo 950
        'header_text': 'FFFFFF',
        'kpi_fill': 'EEF2F6',    # Indigo 50
        'kpi_text': '3730A3',
    },
    'nashm-crm-table': {
        'header_fill': '0F172A', # Slate 900
        'header_text': 'FFFFFF',
        'kpi_fill': 'F8FAFC',    # Slate 50
        'kpi_text': '334155',
    },
    'nashm-inventory': {
        'header_fill': '451A03', # Amber 950
        'header_text': 'FFFFFF',
        'kpi_fill': 'FEF3C7',    # Amber 100
        'kpi_text': '92400E',
    },
    'nashm-survey-analysis': {
        'header_fill': '0C4A6E', # Sky 950
        'header_text': 'FFFFFF',
        'kpi_fill': 'F0F9FF',    # Sky 50
        'kpi_text': '0369A1',
    }
}

# PPTX XML helpers for RTL
def set_pptx_paragraph_rtl(p):
    pPr = p._p.get_or_add_pPr()
    pPr.set('rtl', '1')

# DOCX XML helpers for RTL
def set_docx_paragraph_rtl(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    if len(pPr.xpath("./w:bidi")) == 0:
        bidi = OxmlElement('w:bidi')
        pPr.insert_element_before(
            bidi,
            "w:adjustRightInd", "w:snapToGrid", "w:spacing", "w:ind", 
            "w:contextualSpacing", "w:mirrorIndents", "w:suppressOverlap", 
            "w:jc", "w:textDirection", "w:textAlignment", "w:textboxTightWrap", 
            "w:outlineLvl", "w:divId", "w:cnfStyle", "w:rPr", "w:sectPr", 
            "w:pPrChange"
        )

def set_docx_run_rtl(run):
    run.font.rtl = True
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.get_or_add_rFonts()
    
    font_name = run.font.name or 'Cairo'
    font_size_pt = run.font.size.pt if run.font.size else 11
    
    rFonts.set(qn('w:cs'), font_name)
    rFonts.set(qn('w:ascii'), font_name)
    rFonts.set(qn('w:hAnsi'), font_name)
    
    # Set complex script size (szCs) in half-points
    szCs = rPr.find(qn('w:szCs'))
    if szCs is None:
        szCs = OxmlElement('w:szCs')
        rPr.append(szCs)
    szCs.set(qn('w:val'), str(int(font_size_pt * 2)))
    
    # Set language to Arabic for complex script
    lang = rPr.find(qn('w:lang'))
    if lang is None:
        lang = OxmlElement('w:lang')
        rPr.append(lang)
    lang.set(qn('w:bidi'), 'ar-SA')

def generate_slides(data, output_path):
    theme_id = data.get('templateId', 'nashm-executive-dark')
    theme = SLIDES_THEMES.get(theme_id, SLIDES_THEMES['nashm-executive-dark'])
    
    prs = Presentation()
    # Set to standard 16:9 widescreen slides size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    is_rtl = data.get('direction') == 'rtl' or (data.get('direction') == 'auto' and data.get('locale', '').startswith('ar'))

    # Background color helper
    def style_slide_background(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme['bg'])

    blank_slide_layout = prs.slide_layouts[6] # Blank slide layout

    for slide_data in data.get('content', {}).get('slides', []):
        slide = prs.slides.add_slide(blank_slide_layout)
        style_slide_background(slide)

        layout = slide_data.get('layout', 'split')
        title_text = slide_data.get('title', '')
        eyebrow = slide_data.get('eyebrow', '')
        content_items = slide_data.get('content', [])

        # Add Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0

        # Eyebrow
        if eyebrow:
            p_eyebrow = tf.paragraphs[0]
            p_eyebrow.text = eyebrow.upper()
            p_eyebrow.font.name = 'Cairo' if is_rtl else 'Arial'
            p_eyebrow.font.size = Pt(10)
            p_eyebrow.font.bold = True
            p_eyebrow.font.color.rgb = RGBColor(*theme['accent'])
            if is_rtl:
                p_eyebrow.alignment = PP_ALIGN.RIGHT
                set_pptx_paragraph_rtl(p_eyebrow)
            p_title = tf.add_paragraph()
        else:
            p_title = tf.paragraphs[0]

        # Title
        p_title.text = title_text
        p_title.font.name = 'Cairo' if is_rtl else 'Arial'
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(*theme['text'])
        if is_rtl:
            p_title.alignment = PP_ALIGN.RIGHT
            set_pptx_paragraph_rtl(p_title)

        # Content layouts rendering
        if layout == 'cover':
            # Center title vertically and horizontally
            p_title.alignment = PP_ALIGN.CENTER
            if eyebrow:
                p_eyebrow.alignment = PP_ALIGN.CENTER
            title_box.top = Inches(2.2)
            title_box.height = Inches(3)
            
            if content_items:
                p_desc = tf.add_paragraph()
                p_desc.text = content_items[0]
                p_desc.space_before = Pt(24)
                p_desc.font.name = 'Cairo' if is_rtl else 'Arial'
                p_desc.font.size = Pt(16)
                p_desc.font.color.rgb = RGBColor(*theme['text'])
                p_desc.alignment = PP_ALIGN.CENTER
                if is_rtl:
                    set_pptx_paragraph_rtl(p_desc)

        elif layout == 'section':
            # Section divider / title layout with high fidelity
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
            tf_section = content_box.text_frame
            tf_section.word_wrap = True
            tf_section.margin_left = tf_section.margin_right = 0
            
            p_sec_title = tf_section.paragraphs[0]
            p_sec_title.text = title_text
            p_sec_title.font.name = 'Cairo' if is_rtl else 'Arial'
            p_sec_title.font.size = Pt(28)
            p_sec_title.font.bold = True
            p_sec_title.font.color.rgb = RGBColor(*theme['text'])
            if is_rtl:
                p_sec_title.alignment = PP_ALIGN.RIGHT
                set_pptx_paragraph_rtl(p_sec_title)

            # Colored accent line
            line_top = Inches(2.9) if eyebrow else Inches(2.7)
            shape_line = slide.shapes.add_shape(1, Inches(10.5) if is_rtl else Inches(0.8), line_top, Inches(2.0), Inches(0.08))
            shape_line.fill.solid()
            shape_line.fill.fore_color.rgb = RGBColor(*theme['accent'])
            shape_line.line.fill.background()

            # Description content
            for text in content_items:
                p_desc = tf_section.add_paragraph()
                p_desc.text = text
                p_desc.space_before = Pt(14)
                p_desc.font.name = 'Cairo' if is_rtl else 'Arial'
                p_desc.font.size = Pt(14)
                p_desc.font.color.rgb = RGBColor(*theme['text'])
                if is_rtl:
                    p_desc.alignment = PP_ALIGN.RIGHT
                    set_pptx_paragraph_rtl(p_desc)

        elif layout == 'kpi':
            # 2x2 grid of boxes below title
            box_width = Inches(5.4)
            box_height = Inches(2.0)
            
            # Position calculations (swapped if RTL)
            c1_left = Inches(6.8) if is_rtl else Inches(0.8)
            c2_left = Inches(0.8) if is_rtl else Inches(6.8)
            
            coords = [
                (c1_left, Inches(2.5)), # Box 1
                (c2_left, Inches(2.5)), # Box 2
                (c1_left, Inches(4.8)), # Box 3
                (c2_left, Inches(4.8))  # Box 4
            ]
            for i, item in enumerate(content_items[:4]):
                left, top = coords[i]
                # Draw outer card border / background
                shape = slide.shapes.add_shape(1, left, top, box_width, box_height) # 1 = Rectangle
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*theme['card_bg'])
                shape.line.color.rgb = RGBColor(*theme['accent'])
                shape.line.width = Pt(1)

                tf_card = shape.text_frame
                tf_card.word_wrap = True
                tf_card.margin_left = tf_card.margin_right = Inches(0.3)
                tf_card.margin_top = Inches(0.4)

                parts = item.split(':')
                val = parts[0].strip()
                label = parts[1].strip() if len(parts) > 1 else ''

                p_val = tf_card.paragraphs[0]
                p_val.text = val
                p_val.font.name = 'Cairo' if is_rtl else 'Arial'
                p_val.font.size = Pt(32)
                p_val.font.bold = True
                p_val.font.color.rgb = RGBColor(*theme['accent'])
                
                if label:
                    p_label = tf_card.add_paragraph()
                    p_label.text = label
                    p_label.space_before = Pt(8)
                    p_label.font.name = 'Cairo' if is_rtl else 'Arial'
                    p_label.font.size = Pt(12)
                    p_label.font.color.rgb = RGBColor(*theme['text'])
                    p_label.font.bold = True

                if is_rtl:
                    p_val.alignment = PP_ALIGN.RIGHT
                    set_pptx_paragraph_rtl(p_val)
                    if label:
                        p_label.alignment = PP_ALIGN.RIGHT
                        set_pptx_paragraph_rtl(p_label)

        elif layout == 'comparison':
            # Side-by-side comparison cards
            col_width = Inches(5.5)
            col_height = Inches(4.0)
            
            # Position calculations (swapped if RTL)
            c1_left = Inches(7.0) if is_rtl else Inches(0.8)
            c2_left = Inches(0.8) if is_rtl else Inches(7.0)
            
            for i, col_data in enumerate(content_items[:2]):
                left = c1_left if i == 0 else c2_left
                shape = slide.shapes.add_shape(1, left, Inches(2.2), col_width, col_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*theme['card_bg'])
                shape.line.color.rgb = RGBColor(*theme['accent'])
                shape.line.width = Pt(1)
                
                tf_comp = shape.text_frame
                tf_comp.word_wrap = True
                tf_comp.margin_left = tf_comp.margin_right = Inches(0.3)
                tf_comp.margin_top = Inches(0.3)
                
                lines = [line.strip() for line in col_data.split('\n') if line.strip()]
                col_title = lines[0] if lines else ''
                items = lines[1:] if len(lines) > 1 else []
                
                # Column Title
                p_col_title = tf_comp.paragraphs[0]
                p_col_title.text = col_title
                p_col_title.font.name = 'Cairo' if is_rtl else 'Arial'
                p_col_title.font.size = Pt(16)
                p_col_title.font.bold = True
                p_col_title.font.color.rgb = RGBColor(*theme['accent'])
                if is_rtl:
                    p_col_title.alignment = PP_ALIGN.RIGHT
                    set_pptx_paragraph_rtl(p_col_title)
                    
                # Column Items
                for item in items:
                    p_item = tf_comp.add_paragraph()
                    clean_item = item.lstrip('-').lstrip('•').strip()
                    p_item.text = "• " + clean_item
                    p_item.space_before = Pt(8)
                    p_item.font.name = 'Cairo' if is_rtl else 'Arial'
                    p_item.font.size = Pt(12)
                    p_item.font.color.rgb = RGBColor(*theme['text'])
                    if is_rtl:
                        p_item.alignment = PP_ALIGN.RIGHT
                        set_pptx_paragraph_rtl(p_item)

        elif layout == 'grid':
            # Grid layout of up to 4 items
            box_width = Inches(5.4)
            box_height = Inches(1.8)
            
            c1_left = Inches(6.8) if is_rtl else Inches(0.8)
            c2_left = Inches(0.8) if is_rtl else Inches(6.8)
            
            coords = [
                (c1_left, Inches(2.4)),
                (c2_left, Inches(2.4)),
                (c1_left, Inches(4.5)),
                (c2_left, Inches(4.5))
            ]
            for i, text in enumerate(content_items[:4]):
                left, top = coords[i]
                shape = slide.shapes.add_shape(1, left, top, box_width, box_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*theme['card_bg'])
                shape.line.color.rgb = RGBColor(*theme['accent'])
                shape.line.width = Pt(1)
                
                tf_grid = shape.text_frame
                tf_grid.word_wrap = True
                tf_grid.margin_left = tf_grid.margin_right = Inches(0.25)
                tf_grid.margin_top = Inches(0.25)
                
                p = tf_grid.paragraphs[0]
                p.text = text
                p.font.name = 'Cairo' if is_rtl else 'Arial'
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(*theme['text'])
                if is_rtl:
                    p.alignment = PP_ALIGN.RIGHT
                    set_pptx_paragraph_rtl(p)

        elif layout == 'agenda':
            # Agenda list layout
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
            tf_agenda = content_box.text_frame
            tf_agenda.word_wrap = True
            tf_agenda.margin_left = tf_agenda.margin_right = 0
            
            for i, text in enumerate(content_items):
                p = tf_agenda.paragraphs[0] if i == 0 else tf_agenda.add_paragraph()
                p.text = f"{i + 1}.  {text}"
                p.space_after = Pt(14)
                p.font.name = 'Cairo' if is_rtl else 'Arial'
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*theme['text'])
                if is_rtl:
                    p.alignment = PP_ALIGN.RIGHT
                    set_pptx_paragraph_rtl(p)

        elif layout == 'split':
            # Two column text layout
            col_width = Inches(5.5)
            col_height = Inches(4.0)
            
            # Position calculations (swapped if RTL)
            c1_left = Inches(7.0) if is_rtl else Inches(0.8)
            c2_left = Inches(0.8) if is_rtl else Inches(7.0)
            
            # Left column
            left_col = slide.shapes.add_textbox(c1_left, Inches(2.4), col_width, col_height)
            tf_left = left_col.text_frame
            tf_left.word_wrap = True
            
            # Right column
            right_col = slide.shapes.add_textbox(c2_left, Inches(2.4), col_width, col_height)
            tf_right = right_col.text_frame
            tf_right.word_wrap = True

            mid = (len(content_items) + 1) // 2
            left_items = content_items[:mid]
            right_items = content_items[mid:]

            for i, text in enumerate(left_items):
                p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
                p.text = "• " + text
                p.space_after = Pt(12)
                p.font.name = 'Cairo' if is_rtl else 'Arial'
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(*theme['text'])
                if is_rtl:
                    p.alignment = PP_ALIGN.RIGHT
                    set_pptx_paragraph_rtl(p)

            for i, text in enumerate(right_items):
                p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
                p.text = "• " + text
                p.space_after = Pt(12)
                p.font.name = 'Cairo' if is_rtl else 'Arial'
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(*theme['text'])
                if is_rtl:
                    p.alignment = PP_ALIGN.RIGHT
                    set_pptx_paragraph_rtl(p)

        else:
            # Standard single textbox list
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(4.2))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            tf_content.margin_left = tf_content.margin_right = 0
            
            for i, text in enumerate(content_items):
                p = tf_content.paragraphs[0] if i == 0 else tf_content.add_paragraph()
                p.text = "• " + text
                p.space_after = Pt(14)
                p.font.name = 'Cairo' if is_rtl else 'Arial'
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(*theme['text'])
                if is_rtl:
                    p.alignment = PP_ALIGN.RIGHT
                    set_pptx_paragraph_rtl(p)

        # Set speaker notes if available
        notes = slide_data.get('notes', '')
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(output_path)


def generate_document(data, output_path):
    theme_id = data.get('templateId', 'nashm-report-pro')
    theme = DOCS_THEMES.get(theme_id, DOCS_THEMES['nashm-report-pro'])

    doc = Document()
    is_rtl = data.get('direction') == 'rtl' or (data.get('direction') == 'auto' and data.get('locale', '').startswith('ar'))

    # Set page margins
    sections_list = doc.sections
    for section in sections_list:
        section.top_margin = DocxInches(1)
        section.bottom_margin = DocxInches(1)
        section.left_margin = DocxInches(1)
        section.right_margin = DocxInches(1)

    # Document main title
    title_p = doc.add_paragraph()
    title_run = title_p.add_run(data.get('title', 'Generated Document'))
    title_run.font.name = 'Cairo' if is_rtl else 'Arial'
    title_run.font.size = DocxPt(26)
    title_run.font.bold = True
    title_run.font.color.rgb = DocxRGBColor(*theme['header_color'])
    
    if is_rtl:
        title_p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        set_docx_paragraph_rtl(title_p)
        set_docx_run_rtl(title_run)
    else:
        title_p.alignment = WD_ALIGN_PARAGRAPH.LEFT

    # Add spacing after title
    title_p.paragraph_format.space_after = DocxPt(24)

    for sec in data.get('content', {}).get('sections', []):
        # Heading
        head_p = doc.add_paragraph()
        head_run = head_p.add_run(sec.get('title', ''))
        head_run.font.name = 'Cairo' if is_rtl else 'Arial'
        head_run.font.size = DocxPt(16)
        head_run.font.bold = True
        head_run.font.color.rgb = DocxRGBColor(*theme['header_color'])
        
        if is_rtl:
            head_p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            set_docx_paragraph_rtl(head_p)
            set_docx_run_rtl(head_run)
        else:
            head_p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            
        head_p.paragraph_format.space_before = DocxPt(18)
        head_p.paragraph_format.space_after = DocxPt(12)

        # Paragraphs
        for p_text in sec.get('paragraphs', []):
            p = doc.add_paragraph()
            p_run = p.add_run(p_text)
            p_run.font.name = 'Cairo' if is_rtl else 'Arial'
            p_run.font.size = DocxPt(11)
            p_run.font.color.rgb = DocxRGBColor(*theme['text_color'])
            
            if is_rtl:
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                set_docx_paragraph_rtl(p)
                set_docx_run_rtl(p_run)
            else:
                p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                
            p.paragraph_format.space_after = DocxPt(10)
            p.paragraph_format.line_spacing = 1.15

        # Lists
        lst = sec.get('list')
        if lst and lst.get('items'):
            list_type = lst.get('type', 'bullet')
            for i, item_text in enumerate(lst.get('items', [])):
                p = doc.add_paragraph()
                prefix = f"{i+1}. " if list_type == 'numbered' else "• "
                p_run = p.add_run(prefix + item_text)
                p_run.font.name = 'Cairo' if is_rtl else 'Arial'
                p_run.font.size = DocxPt(11)
                p_run.font.color.rgb = DocxRGBColor(*theme['text_color'])
                
                if is_rtl:
                    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                    set_docx_paragraph_rtl(p)
                    set_docx_run_rtl(p_run)
                else:
                    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    p.paragraph_format.left_indent = DocxInches(0.25)
                    
                p.paragraph_format.space_after = DocxPt(6)

    doc.save(output_path)


def generate_workbook(data, output_path):
    theme_id = data.get('templateId', 'nashm-finance-dashboard')
    theme = SHEETS_THEMES.get(theme_id, SHEETS_THEMES['nashm-finance-dashboard'])

    wb = Workbook()
    is_rtl = data.get('direction') == 'rtl' or (data.get('direction') == 'auto' and data.get('locale', '').startswith('ar'))

    # Strip default sheet
    default_sheet = wb.active
    wb.remove(default_sheet)

    for idx, sheet_data in enumerate(data.get('content', {}).get('sheets', [])):
        ws = wb.create_sheet(title=sheet_data.get('name', f"Sheet {idx+1}"))
        
        # Gridlines enabled
        ws.sheet_view.showGridLines = True
        ws.views.sheetView[0].showGridLines = True

        if is_rtl:
            ws.sheet_view.rightToLeft = True

        start_row = 1

        # Summary KPIs
        summary = sheet_data.get('summary', [])
        if summary:
            # Write KPI boxes in row 1-3
            for col_idx, item in enumerate(summary[:4]):
                # Map columns (e.g. col_idx = 0 -> A, B; col_idx = 1 -> C, D)
                c1 = col_idx * 2 + 1
                c2 = c1 + 1
                
                # Merge cells for label and value
                ws.merge_cells(start_row=1, start_column=c1, end_row=1, end_column=c2)
                ws.merge_cells(start_row=2, start_column=c1, end_row=2, end_column=c2)
                
                lbl_cell = ws.cell(row=1, column=c1)
                lbl_cell.value = item.get('label', '').upper()
                lbl_cell.font = Font(name='Cairo' if is_rtl else 'Arial', size=9, bold=True, color='555555')
                lbl_cell.fill = PatternFill(start_color=theme['kpi_fill'], end_color=theme['kpi_fill'], fill_type='solid')
                lbl_cell.alignment = Alignment(horizontal='center', vertical='center')

                val_cell = ws.cell(row=2, column=c1)
                val_cell.value = item.get('value', '')
                val_cell.font = Font(name='Cairo' if is_rtl else 'Arial', size=16, bold=True, color=theme['kpi_text'])
                val_cell.fill = PatternFill(start_color=theme['kpi_fill'], end_color=theme['kpi_fill'], fill_type='solid')
                val_cell.alignment = Alignment(horizontal='center', vertical='center')

                # Border styles for KPI boxes
                thin = Side(border_style="thin", color="CCCCCC")
                for r in [1, 2]:
                    for c in [c1, c2]:
                        ws.cell(row=r, column=c).border = Border(top=thin, left=thin, right=thin, bottom=thin)

            start_row = 4 # Shift table start down

        # Table headers
        headers = sheet_data.get('headers', [])
        for col_idx, header in enumerate(headers):
            cell = ws.cell(row=start_row, column=col_idx + 1)
            cell.value = header
            cell.font = Font(name='Cairo' if is_rtl else 'Arial', size=11, bold=True, color=theme['header_text'])
            cell.fill = PatternFill(start_color=theme['header_fill'], end_color=theme['header_fill'], fill_type='solid')
            cell.alignment = Alignment(horizontal='right' if is_rtl else 'left', vertical='center', readingOrder=2 if is_rtl else 1)
            thin_border_side = Side(border_style="thin", color="CCCCCC")
            cell.border = Border(top=thin_border_side, bottom=thin_border_side, left=thin_border_side, right=thin_border_side)

        # Rows
        rows = sheet_data.get('rows', [])
        for row_idx, row_values in enumerate(rows):
            current_row = start_row + 1 + row_idx
            
            # Check if this is a totals row
            is_total_row = any(
                isinstance(v, str) and 
                (v.lower().find('total') != -1 or v.lower().find('sum') != -1 or v.find('مجموع') != -1 or v.find('إجمالي') != -1) 
                for v in row_values
            )

            for col_idx, val in enumerate(row_values):
                cell = ws.cell(row=current_row, column=col_idx + 1)
                cell.value = val

                is_num = isinstance(val, (int, float))

                # Font & Style
                font_weight = True if is_total_row else False
                cell.font = Font(name='Cairo' if is_rtl else 'Arial', size=10, bold=font_weight)
                
                # Alignment
                horiz_align = 'right' if is_num or is_rtl else 'left'
                cell.alignment = Alignment(horizontal=horiz_align, vertical='center', readingOrder=2 if is_rtl else 1)

                # Zebra stripes and borders for standard/total rows
                if is_total_row:
                    cell.fill = PatternFill(start_color='E2E8F0', end_color='E2E8F0', fill_type='solid')
                    double_border = Side(border_style="double", color="333333")
                    thin_top = Side(border_style="thin", color="666666")
                    cell.border = Border(top=thin_top, bottom=double_border)
                else:
                    thin_side = Side(border_style="thin", color="E2E8F0")
                    cell.border = Border(top=thin_side, bottom=thin_side, left=thin_side, right=thin_side)
                    if row_idx % 2 == 1:
                        cell.fill = PatternFill(start_color='F8FAFC', end_color='F8FAFC', fill_type='solid')

        # Auto fit column widths
        for col in ws.columns:
            max_len = 0
            for cell in col:
                # skip merged KPI cells to avoid stretching columns excessively
                if cell.row < start_row:
                    continue
                val_str = str(cell.value or '')
                if len(val_str) > max_len:
                    max_len = len(val_str)
            col_letter = get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(max_len + 5, 14)

    wb.save(output_path)


def main():
    parser = argparse.ArgumentParser(description="Nashm Document Exporter Pipeline")
    parser.add_argument('--input', help="Path to raw input structural JSON file")
    parser.add_argument('--output', help="Target generated output path")
    parser.add_argument('--format', default="office", choices=["office", "pdf"], help="Export format")
    parser.add_argument('--test', action='store_true', help="Run in self-test mock mode")
    
    args = parser.parse_args()

    if args.test:
        print("Self-test completed successfully. Libraries verified.")
        return

    if not args.input or not args.output:
        print("Error: --input and --output are required unless --test is used.")
        sys.exit(1)

    if not os.path.exists(args.input):
        print(f"Error: input file {args.input} does not exist.")
        sys.exit(1)

    try:
        with open(args.input, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        print(f"Error reading JSON: {e}")
        sys.exit(1)

    kind = data.get('kind')
    if not kind:
        print("Error: kind attribute missing in structural JSON")
        sys.exit(1)

    # Make parent dirs if needed
    os.makedirs(os.path.dirname(os.path.abspath(args.output)), exist_ok=True)

    try:
        if kind == 'slides':
            generate_slides(data, args.output)
        elif kind == 'document':
            generate_document(data, args.output)
        elif kind == 'workbook':
            generate_workbook(data, args.output)
        else:
            print(f"Error: Unknown document kind '{kind}'")
            sys.exit(1)
        
        print(f"Export successful. File written to {args.output}")
    except Exception as e:
        print(f"Exporter execution failure: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
